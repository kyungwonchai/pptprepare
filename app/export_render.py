"""선택한 슬라이드 URL을 Playwright로 캡처해 PNG 바이트 목록으로 반환."""

from __future__ import annotations

import io
import os
import tempfile
import time


def slide_path_for_export(n: int) -> str:
    """브라우저에서 열 경로 (슬라이드 번호 → URL path + export 쿼리)."""
    if n == 1:
        return "/?export=1"
    return f"/{n}?export=1"


def capture_slide_frames(
    base_url: str,
    slide_numbers: list[int],
    theme: str,
) -> list[bytes]:
    """
    base_url: 예) http://127.0.0.1:54321/
    theme: 'dark' | 'light'
    """
    try:
        from playwright.sync_api import sync_playwright
    except ImportError as e:
        raise RuntimeError(
            "Playwright가 설치되어 있지 않습니다. "
            "`pip install playwright` 후 `playwright install chromium` 을 실행하세요."
        ) from e

    if theme not in ("dark", "light"):
        theme = "dark"

    root = base_url.rstrip("/")
    out: list[bytes] = []

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        try:
            context = browser.new_context(
                viewport={"width": 1280, "height": 720},
                device_scale_factor=2,
            )
            page = context.new_page()
            for n in slide_numbers:
                path = slide_path_for_export(n)
                url = root + path
                page.goto(url, wait_until="domcontentloaded", timeout=60_000)
                page.wait_for_selector(".slide-frame", timeout=30_000)
                page.evaluate(
                    """(t) => {
                      document.documentElement.setAttribute('data-theme', t);
                      document.body.classList.add('export-capture');
                    }""",
                    theme,
                )
                time.sleep(0.35)
                frame = page.locator(".slide-frame").first
                png = frame.screenshot(type="png")
                out.append(png)
            context.close()
        finally:
            browser.close()

    return out


def _blank_layout(prs):
    for layout in prs.slide_layouts:
        name = (layout.name or "").lower()
        if "blank" in name:
            return layout
    if len(prs.slide_layouts) > 6:
        return prs.slide_layouts[6]
    return prs.slide_layouts[0]


def build_pptx_from_images(images: list[bytes]) -> bytes:
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError as e:
        raise RuntimeError("`pip install python-pptx` 가 필요합니다.") from e

    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(5143500)
    layout = _blank_layout(prs)

    for raw in images:
        slide = prs.slides.add_slide(layout)
        stream = io.BytesIO(raw)
        stream.seek(0)
        slide.shapes.add_picture(stream, 0, 0, width=prs.slide_width, height=prs.slide_height)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def build_pdf_from_images(images: list[bytes]) -> bytes:
    try:
        import img2pdf
    except ImportError as e:
        raise RuntimeError("`pip install img2pdf` 가 필요합니다.") from e

    with tempfile.TemporaryDirectory() as tmp:
        paths: list[str] = []
        for i, raw in enumerate(images):
            p = os.path.join(tmp, f"s{i}.png")
            with open(p, "wb") as f:
                f.write(raw)
            paths.append(p)
        return img2pdf.convert(paths)
